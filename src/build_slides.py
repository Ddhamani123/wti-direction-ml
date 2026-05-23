"""Build the 4-slide VIP presentation deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
from pathlib import Path
import copy

# ── helpers ──────────────────────────────────────────────────────────────────

DATA = Path("data")

# Colour palette
NAVY   = RGBColor(0x1A, 0x23, 0x7E)   # dark blue  – headings
TEAL   = RGBColor(0x00, 0x89, 0x7B)   # teal       – accents / step boxes
AMBER  = RGBColor(0xF9, 0xA8, 0x25)   # amber      – arrow / highlight
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)   # light grey – slide bg tint
DGRAY  = RGBColor(0x42, 0x42, 0x42)   # dark grey  – body text
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
RED    = RGBColor(0xC6, 0x28, 0x28)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]      # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor = LGRAY):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h,
        text="", font_size=14, bold=False, color=DGRAY,
        bg_color=None, align=PP_ALIGN.LEFT,
        font_color=None, border_color=None):
    """Add a rounded-corner text box."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    # background fill
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color

    # border
    if border_color:
        ln = txBox.line
        ln.color.rgb = border_color
        ln.width = Pt(1)

    # text
    if text:
        from pptx.util import Pt as _Pt
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = _Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = font_color or color
    return txBox


def title_bar(slide, text, subtitle=""):
    """Full-width navy title bar at top."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(13.33), Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = "  " + subtitle
        r2.font.size = Pt(13)
        r2.font.bold = False
        r2.font.color.rgb = RGBColor(0xB0, 0xBE, 0xC5)


def step_box(slide, l, t, w, h, number, title, body, bg_c=TEAL):
    """Numbered pipeline step box."""
    # coloured box
    rect = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_c
    rect.line.fill.background()

    # number circle (fake with a small box)
    num_rect = slide.shapes.add_shape(1,
        Inches(l + 0.08), Inches(t + 0.08),
        Inches(0.45), Inches(0.45))
    num_rect.fill.solid()
    num_rect.fill.fore_color.rgb = WHITE
    num_rect.line.fill.background()
    tf_n = num_rect.text_frame
    pn = tf_n.paragraphs[0]
    pn.alignment = PP_ALIGN.CENTER
    rn = pn.add_run()
    rn.text = str(number)
    rn.font.size = Pt(14)
    rn.font.bold = True
    rn.font.color.rgb = bg_c

    # title
    tb = slide.shapes.add_textbox(
        Inches(l + 0.62), Inches(t + 0.10), Inches(w - 0.72), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # body
    tb2 = slide.shapes.add_textbox(
        Inches(l + 0.15), Inches(t + 0.52), Inches(w - 0.25), Inches(h - 0.60))
    tb2.text_frame.word_wrap = True
    p2 = tb2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(11)
    r2.font.color.rgb = WHITE


def arrow(slide, l, t, length=0.35):
    """Downward arrow between step boxes."""
    arr = slide.shapes.add_shape(
        13,  # MSO_SHAPE_TYPE.DOWN_ARROW  (enum val 13)
        Inches(l), Inches(t), Inches(0.30), Inches(length)
    )
    arr.fill.solid()
    arr.fill.fore_color.rgb = AMBER
    arr.line.fill.background()


def img(slide, path, l, t, w, h):
    slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))


def section_label(slide, text, l, t, w=3.5):
    """Small teal label above a section."""
    b = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.28))
    b.fill.solid()
    b.fill.fore_color.rgb = TEAL
    b.line.fill.background()
    tf = b.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE


def divider(slide, t):
    line = slide.shapes.add_shape(1, Inches(0.3), Inches(t), Inches(12.73), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    line.line.fill.background()


# ── Slide 1: Sentiment Pipeline ───────────────────────────────────────────────

def slide1(prs):
    sl = blank_slide(prs)
    bg(sl, WHITE)
    title_bar(sl, "News Sentiment Pipeline",
              "Yahoo Finance  ·  FinBERT  ·  21 Energy Tickers")

    # ── Left: pipeline steps ──
    section_label(sl, "DATA COLLECTION PIPELINE", 0.35, 1.30, 4.5)

    step_box(sl, 0.35, 1.65, 4.5, 1.15, 1, "FETCH  —  yfinance",
             "21 energy tickers (WTI/Brent futures, US & intl majors, oilfield\n"
             "services, ETFs, natural gas)  ·  deduplicated by headline title",
             bg_c=RGBColor(0x1E, 0x88, 0xE5))
    arrow(sl, 2.27, 2.82)
    step_box(sl, 0.35, 3.20, 4.5, 1.15, 2, "SCORE  —  ProsusAI/FinBERT (GPU)",
             "BERT fine-tuned on financial text  ·  batch=32  ·  labels: positive /\n"
             "negative / neutral  →  signed score: +conf / −conf / 0",
             bg_c=TEAL)
    arrow(sl, 2.27, 4.37)
    step_box(sl, 0.35, 4.75, 4.5, 1.15, 3, "AGGREGATE  —  7 daily features",
             "avg_sentiment · positive_pct · negative_pct · neutral_pct\n"
             "article_count · sentiment_5d_ma · sentiment_20d_ma",
             bg_c=RGBColor(0x6A, 0x1B, 0x9A))

    # coverage note
    box(sl, 0.35, 6.05, 4.5, 0.60,
        "Coverage: ~10 months  ·  170 unique headlines  ·  cached to CSV",
        font_size=11, color=DGRAY, bg_color=RGBColor(0xE8, 0xF5, 0xE9),
        border_color=GREEN)

    # ── Vertical divider ──
    div = sl.shapes.add_shape(1, Inches(5.05), Inches(1.25), Inches(0.02), Inches(5.8))
    div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
    div.line.fill.background()

    # ── Right top: article count ──
    section_label(sl, "MONTHLY ARTICLE COVERAGE", 5.25, 1.30, 3.8)
    img(sl, DATA/"fig_article_count.png", 5.20, 1.62, 3.9, 2.55)

    # ── Right bottom: label distribution ──
    section_label(sl, "FINBERT LABEL DISTRIBUTION", 9.25, 1.30, 3.75)
    img(sl, DATA/"fig_sentiment_vs_price.png", 9.20, 1.62, 3.9, 2.55)   # placeholder; replace below

    # use the label dist figure
    # (fig_sentiment_vs_price is the daily chart — use article_count + wordclouds for slide 1 right)
    # We'll put label dist at bottom right instead
    section_label(sl, "FINBERT LABEL DISTRIBUTION", 9.25, 4.30, 3.75)
    img(sl, DATA/"fig_wordclouds.png", 5.20, 4.30, 7.8, 2.85)

    return sl


# ── Slide 2: Sentiment Analysis ───────────────────────────────────────────────

def slide2(prs):
    sl = blank_slide(prs)
    bg(sl, WHITE)
    title_bar(sl, "Sentiment Analysis & WTI Integration",
              "Preprocessing  ·  Feature Engineering  ·  Proof of Concept")

    # hero: daily sentiment + MAs
    section_label(sl, "DAILY SENTIMENT  +  5d / 20d MOVING AVERAGES", 0.3, 1.25, 6.5)
    img(sl, DATA/"fig_sentiment_vs_price.png", 0.3, 1.58, 7.8, 3.0)

    # word clouds
    section_label(sl, "POSITIVE vs NEGATIVE WORD CLOUDS", 8.35, 1.25, 4.65)
    img(sl, DATA/"fig_wordclouds.png", 8.30, 1.58, 4.70, 3.0)

    divider(sl, 4.75)

    # integration design bullets
    section_label(sl, "PROOF OF CONCEPT  —  WTI FEATURE MATRIX INTEGRATION", 0.3, 4.88, 7.0)

    bullets = [
        ("Base features", "31 FRED-MD macro vars  +  30 lagged WTI closing prices"),
        ("Sentiment added", "avg_sentiment · positive_pct · negative_pct · sentiment_5d_ma"),
        ("No data leakage", "All sentiment features lagged 1 trading day before merging"),
        ("Gap handling", "Forward-fill up to 5 business days on days without articles"),
        ("Overlap window", "133 trading days demonstrated (May – Oct 2025)"),
    ]
    y = 5.20
    for label, detail in bullets:
        box(sl, 0.40, y, 2.10, 0.32, label,
            font_size=10, bold=True, font_color=TEAL,
            bg_color=RGBColor(0xE0, 0xF2, 0xF1))
        box(sl, 2.55, y, 9.80, 0.32, detail,
            font_size=10, color=DGRAY)
        y += 0.38

    return sl


# ── Slide 3: Strategy Explanation ────────────────────────────────────────────

def slide3(prs):
    sl = blank_slide(prs)
    bg(sl, WHITE)
    title_bar(sl, "WTI Direction-Based Trading Strategies",
              "Long-Short  ·  Confidence-Filtered  ·  Ensemble")

    # ── Strategy cards ──
    cards = [
        ("Long-Short",         NAVY,                    RGBColor(0xE8,0xEA,0xF6),
         "+1 (Long)\nif model predicts UP\n\n−1 (Short)\nif model predicts DOWN",
         "Trades every day.\nCaptures every signal."),
        ("Confidence-Filtered", RGBColor(0x00,0x69,0x5C), RGBColor(0xE0,0xF2,0xF1),
         "+1 / −1\nonly when\n|prob − 0.5| > θ\n\n0 (Flat)\nwhen uncertain",
         "Avoids low-conviction\ntrades.\nThreshold θ tuned on\nvalidation set."),
        ("Ensemble",            RGBColor(0x6A,0x1B,0x9A), RGBColor(0xF3,0xE5,0xF5),
         "Average XGBoost &\nMLP+EN probabilities\n→ apply ensemble θ",
         "Combines both models.\nMore robust signal."),
    ]

    x = 0.35
    for name, hdr_c, bg_c, position_text, edge_text in cards:
        w = 4.15
        # header
        hdr = sl.shapes.add_shape(1, Inches(x), Inches(1.25), Inches(w), Inches(0.45))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = hdr_c; hdr.line.fill.background()
        tf = hdr.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = name; r.font.size = Pt(16); r.font.bold = True
        r.font.color.rgb = WHITE

        # card body
        card = sl.shapes.add_shape(1, Inches(x), Inches(1.70), Inches(w), Inches(3.30))
        card.fill.solid(); card.fill.fore_color.rgb = bg_c
        card.line.color.rgb = hdr_c; card.line.width = Pt(1)

        # "Position" label
        box(sl, x+0.12, 1.80, 1.1, 0.28, "POSITION",
            font_size=9, bold=True, font_color=hdr_c)
        box(sl, x+0.12, 2.12, w-0.24, 1.35, position_text,
            font_size=12, color=DGRAY, align=PP_ALIGN.CENTER)

        # divider
        dl = sl.shapes.add_shape(1, Inches(x+0.15), Inches(3.50), Inches(w-0.30), Inches(0.02))
        dl.fill.solid(); dl.fill.fore_color.rgb = hdr_c; dl.line.fill.background()

        box(sl, x+0.12, 3.55, 1.0, 0.28, "EDGE",
            font_size=9, bold=True, font_color=hdr_c)
        box(sl, x+0.12, 3.85, w-0.24, 1.05, edge_text,
            font_size=11, color=DGRAY)

        x += w + 0.17

    # ── Models used ──
    divider(sl, 5.15)
    section_label(sl, "MODELS", 0.35, 5.28, 1.4)
    box(sl, 1.85, 5.28, 11.0, 0.32,
        "XGBoost  ·  MLP + Elastic Net  —  both output calibrated probabilities used as signals",
        font_size=12, color=DGRAY)

    # ── Key design choices ──
    section_label(sl, "DESIGN CHOICES", 0.35, 5.75, 2.0)
    choices = (
        "Threshold tuned on validation only — no look-ahead bias  ·  "
        "Transaction cost: 0.05 % per trade modeled  ·  "
        "Baseline: Buy & Hold (always long)"
    )
    box(sl, 2.45, 5.75, 10.5, 0.32, choices, font_size=11, color=DGRAY)

    # ── Timeline ──
    section_label(sl, "TEST PERIOD", 0.35, 6.22, 1.7)
    box(sl, 2.15, 6.22, 10.8, 0.32,
        "Jan 2025 – Oct 2025  (out-of-sample, thresholds frozen from validation)",
        font_size=11, color=DGRAY)

    return sl


# ── Slide 4: Strategy Results ─────────────────────────────────────────────────

def slide4(prs):
    sl = blank_slide(prs)
    bg(sl, WHITE)
    title_bar(sl, "Investment Strategy — Results",
              "Test Period: Jan 2025 – Oct 2025  ·  Out-of-Sample")

    # cumulative returns (large, left)
    section_label(sl, "CUMULATIVE RETURNS  —  TEST PERIOD", 0.3, 1.25, 5.5)
    img(sl, DATA/"fig_investment_strategy.png", 0.3, 1.58, 7.55, 3.55)

    # rolling sharpe (right)
    section_label(sl, "ROLLING 60-DAY SHARPE RATIO", 8.05, 1.25, 5.0)
    img(sl, DATA/"fig_rolling_sharpe.png", 8.00, 1.58, 5.05, 3.55)

    divider(sl, 5.28)

    # metrics table
    section_label(sl, "TEST PERIOD SUMMARY", 0.3, 5.40, 2.7)

    headers = ["Strategy", "Total Return", "Sharpe", "Max Drawdown", "Win Rate"]
    rows = [
        ("MLP+EN Long-Short",  "+17.6%", "0.745", "−20.2%", "52.3%", GREEN),
        ("XGBoost Long-Short", "−0.9%",  "0.131", "−29.8%", "49.1%", DGRAY),
        ("Ensemble",           "−9.3%",  "−0.490","−19.8%", "57.4%", DGRAY),
        ("Buy & Hold",         "−14.4%", "−0.397","−27.4%", "49.5%", RED),
    ]
    col_x    = [0.35, 3.05, 5.45, 7.10, 9.55]
    col_w    = [2.55, 2.20, 1.45, 2.30, 2.00]
    row_h    = 0.30

    # header row
    for j, (htext, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        hb = sl.shapes.add_shape(1, Inches(cx), Inches(5.78), Inches(cw), Inches(row_h))
        hb.fill.solid(); hb.fill.fore_color.rgb = NAVY; hb.line.fill.background()
        tf = hb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = htext; r.font.size = Pt(10); r.font.bold = True
        r.font.color.rgb = WHITE

    # data rows
    for i, (strat, ret, sharpe, dd, wr, txt_c) in enumerate(rows):
        ty = 5.78 + row_h + i * row_h
        bg_row = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else WHITE
        vals = [strat, ret, sharpe, dd, wr]
        for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
            rb = sl.shapes.add_shape(1, Inches(cx), Inches(ty), Inches(cw), Inches(row_h))
            rb.fill.solid(); rb.fill.fore_color.rgb = bg_row
            rb.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); rb.line.width = Pt(0.5)
            tf = rb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = val; r.font.size = Pt(10)
            r.font.color.rgb = txt_c if j > 0 else DGRAY
            r.font.bold = (i == 0)   # bold best strategy

    # footnote
    box(sl, 0.35, 7.12, 12.6, 0.25,
        "MLP+EN Long-Short outperforms Buy & Hold by +32 pp  ·  "
        "Transaction costs: 0.05 % per trade modeled in all strategies",
        font_size=9, color=RGBColor(0x75,0x75,0x75), align=PP_ALIGN.CENTER)

    return sl


# ── MAIN ──────────────────────────────────────────────────────────────────────

prs = new_prs()

# Override slide 1 with cleaner layout
sl1 = blank_slide(prs)
bg(sl1, WHITE)
title_bar(sl1, "News Sentiment Pipeline",
          "Data Collection  ·  Yahoo Finance  ·  FinBERT NLP  ·  21 Energy Tickers")

section_label(sl1, "DATA COLLECTION PIPELINE", 0.35, 1.30, 4.5)

step_box(sl1, 0.35, 1.65, 4.5, 1.10, 1, "FETCH  —  yfinance",
         "21 energy tickers (WTI/Brent futures, US & intl majors,\n"
         "oilfield services, ETFs, natural gas)  ·  deduplicated by title",
         bg_c=RGBColor(0x1E, 0x88, 0xE5))
arrow(sl1, 2.27, 2.77)
step_box(sl1, 0.35, 3.15, 4.5, 1.10, 2, "SCORE  —  ProsusAI/FinBERT  (GPU, batch=32)",
         "Fine-tuned BERT on financial text  ·  3 labels: positive / negative /\n"
         "neutral  →  signed score: +confidence / −confidence / 0",
         bg_c=TEAL)
arrow(sl1, 2.27, 4.27)
step_box(sl1, 0.35, 4.65, 4.5, 1.10, 3, "AGGREGATE  —  7 daily features",
         "avg_sentiment · positive_pct · negative_pct · neutral_pct\n"
         "article_count · sentiment_5d_ma · sentiment_20d_ma",
         bg_c=RGBColor(0x6A, 0x1B, 0x9A))

box(sl1, 0.35, 5.90, 4.5, 0.50,
    "Coverage: ~10 months  ·  170 unique headlines  ·  results cached to CSV",
    font_size=11, color=DGRAY, bg_color=RGBColor(0xE8, 0xF5, 0xE9),
    border_color=GREEN)

# vertical divider
div = sl1.shapes.add_shape(1, Inches(5.05), Inches(1.25), Inches(0.02), Inches(5.8))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
div.line.fill.background()

section_label(sl1, "MONTHLY ARTICLE COVERAGE  (Yahoo Finance)", 5.25, 1.30, 5.0)
img(sl1, DATA/"fig_article_count.png", 5.20, 1.62, 4.90, 2.70)

section_label(sl1, "FINBERT LABEL DISTRIBUTION", 10.30, 1.30, 2.70)
# create a small label dist bar inline using a placeholder text since we don't
# have a standalone label-dist png — embed the word cloud bottom half instead
img(sl1, DATA/"fig_wordclouds.png", 5.20, 4.55, 7.80, 2.00)
section_label(sl1, "POSITIVE vs NEGATIVE WORD CLOUDS  (FinBERT)", 5.25, 4.42, 5.5)


# Slide 2
sl2 = blank_slide(prs)
bg(sl2, WHITE)
title_bar(sl2, "Sentiment Analysis & WTI Integration",
          "Preprocessing  ·  Feature Engineering  ·  Proof of Concept")

section_label(sl2, "DAILY SENTIMENT SCORE  +  5-DAY & 20-DAY MOVING AVERAGES", 0.3, 1.25, 8.0)
img(sl2, DATA/"fig_sentiment_vs_price.png", 0.3, 1.58, 8.70, 3.30)

section_label(sl2, "WORD CLOUDS", 9.25, 1.25, 3.80)
img(sl2, DATA/"fig_wordclouds.png", 9.20, 1.58, 3.85, 3.30)

divider(sl2, 5.05)

section_label(sl2, "PROOF OF CONCEPT  —  INTEGRATION INTO WTI FEATURE MATRIX", 0.3, 5.18, 7.5)

bullets = [
    ("Base features",   "31 FRED-MD macro variables  +  30 lagged WTI closing prices"),
    ("+ Sentiment",     "avg_sentiment · positive_pct · negative_pct · sentiment_5d_ma  (4 new features)"),
    ("No leakage",      "All sentiment features shifted 1 trading day — prediction on day t uses sentiment from t−1"),
    ("Gap filling",     "Forward-fill missing sentiment up to 5 business days on days with no articles"),
    ("Overlap window",  "133 trading days demonstrated (May – Oct 2025)  ·  model re-train ready when price data extended"),
]
y = 5.52
for label, detail in bullets:
    b = sl2.shapes.add_shape(1, Inches(0.38), Inches(y), Inches(2.05), Inches(0.30))
    b.fill.solid(); b.fill.fore_color.rgb = TEAL; b.line.fill.background()
    tf = b.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(9.5); r.font.bold = True
    r.font.color.rgb = WHITE
    box(sl2, 2.50, y, 10.55, 0.30, detail, font_size=10, color=DGRAY)
    y += 0.36


# Slide 3
sl3 = blank_slide(prs)
bg(sl3, WHITE)
title_bar(sl3, "WTI Direction-Based Trading Strategies",
          "Long-Short  ·  Confidence-Filtered  ·  Ensemble  ·  Baseline: Buy & Hold")

cards = [
    ("Long-Short",          NAVY,                     RGBColor(0xE8,0xEA,0xF6),
     "+1 (Long)  if model predicts UP\n−1 (Short)  if model predicts DOWN",
     "Trades every day\nCaptures every signal"),
    ("Confidence-Filtered", RGBColor(0x00,0x69,0x5C), RGBColor(0xE0,0xF2,0xF1),
     "+1 / −1  when  |prob − 0.5| > θ\n\n0 (Flat)  when uncertain",
     "Avoids low-conviction trades\nθ tuned on validation set only"),
    ("Ensemble",            RGBColor(0x6A,0x1B,0x9A), RGBColor(0xF3,0xE5,0xF5),
     "Average XGBoost +\nMLP+EN probabilities\nthen apply ensemble θ",
     "Combines both models\nMore robust signal"),
]

x = 0.35
for name, hdr_c, bg_c, pos_text, edge_text in cards:
    w = 4.15
    hdr = sl3.shapes.add_shape(1, Inches(x), Inches(1.25), Inches(w), Inches(0.48))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = hdr_c; hdr.line.fill.background()
    tf = hdr.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name; r.font.size = Pt(17); r.font.bold = True
    r.font.color.rgb = WHITE

    card = sl3.shapes.add_shape(1, Inches(x), Inches(1.73), Inches(w), Inches(3.60))
    card.fill.solid(); card.fill.fore_color.rgb = bg_c
    card.line.color.rgb = hdr_c; card.line.width = Pt(1.5)

    box(sl3, x+0.15, 1.83, w-0.30, 0.28, "POSITION SIZING", font_size=9,
        bold=True, font_color=hdr_c)
    box(sl3, x+0.15, 2.14, w-0.30, 1.60, pos_text,
        font_size=13, color=DGRAY, align=PP_ALIGN.CENTER)

    dl = sl3.shapes.add_shape(1, Inches(x+0.2), Inches(3.76), Inches(w-0.4), Inches(0.02))
    dl.fill.solid(); dl.fill.fore_color.rgb = hdr_c; dl.line.fill.background()

    box(sl3, x+0.15, 3.83, w-0.30, 0.28, "EDGE", font_size=9,
        bold=True, font_color=hdr_c)
    box(sl3, x+0.15, 4.14, w-0.30, 1.10, edge_text,
        font_size=12, color=DGRAY, align=PP_ALIGN.CENTER)

    x += w + 0.17

divider(sl3, 5.48)

row2 = [
    ("MODELS",          "XGBoost  ·  MLP + Elastic Net  —  both output probability estimates used as trading signals"),
    ("THRESHOLD",       "Confidence threshold θ tuned on validation period only → applied frozen to test period (no look-ahead bias)"),
    ("COSTS",           "Transaction cost: 0.05 % per trade modeled  ·  Buy & Hold = always-long baseline"),
]
y = 5.60
for lbl, detail in row2:
    b = sl3.shapes.add_shape(1, Inches(0.38), Inches(y), Inches(1.7), Inches(0.30))
    b.fill.solid(); b.fill.fore_color.rgb = AMBER; b.line.fill.background()
    tf = b.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = lbl; r.font.size = Pt(9.5); r.font.bold = True
    r.font.color.rgb = NAVY
    box(sl3, 2.18, y, 10.9, 0.30, detail, font_size=10, color=DGRAY)
    y += 0.38


# Slide 4
sl4 = blank_slide(prs)
bg(sl4, WHITE)
title_bar(sl4, "Investment Strategy — Results",
          "Test Period: Jan 2025 – Oct 2025  ·  Strictly Out-of-Sample")

section_label(sl4, "CUMULATIVE RETURNS  —  TEST PERIOD", 0.3, 1.25, 5.5)
img(sl4, DATA/"fig_investment_strategy.png", 0.3, 1.58, 7.55, 3.60)

section_label(sl4, "ROLLING 60-DAY SHARPE RATIO", 8.05, 1.25, 5.0)
img(sl4, DATA/"fig_rolling_sharpe.png", 8.00, 1.58, 5.05, 3.60)

divider(sl4, 5.33)

section_label(sl4, "TEST PERIOD PERFORMANCE SUMMARY", 0.3, 5.44, 3.5)

headers = ["Strategy", "Total Return", "Sharpe Ratio", "Max Drawdown", "Win Rate", "Days Traded"]
rows = [
    ("MLP+EN Long-Short",  "+17.6%", "0.745",  "−20.2%", "52.3%", "216", GREEN),
    ("XGBoost Long-Short", "−0.9%",  "0.131",  "−29.8%", "49.1%", "216", DGRAY),
    ("Ensemble",           "−9.3%",  "−0.490", "−19.8%", "57.4%",  "61", DGRAY),
    ("Buy & Hold",         "−14.4%", "−0.397", "−27.4%", "49.5%", "216", RED),
]
col_x = [0.35, 3.00, 4.95, 6.85, 8.75, 10.55]
col_w = [2.50, 1.80, 1.75, 1.75, 1.65,  2.20]
row_h = 0.295

for j, (htext, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    hb = sl4.shapes.add_shape(1, Inches(cx), Inches(5.82), Inches(cw), Inches(row_h))
    hb.fill.solid(); hb.fill.fore_color.rgb = NAVY; hb.line.fill.background()
    tf = hb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = htext; r.font.size = Pt(9.5); r.font.bold = True
    r.font.color.rgb = WHITE

for i, (strat, ret, sharpe, dd, wr, dt, txt_c) in enumerate(rows):
    ty = 5.82 + row_h + i * row_h
    bg_row = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else WHITE
    vals = [strat, ret, sharpe, dd, wr, dt]
    for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
        rb = sl4.shapes.add_shape(1, Inches(cx), Inches(ty), Inches(cw), Inches(row_h))
        rb.fill.solid(); rb.fill.fore_color.rgb = bg_row
        rb.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); rb.line.width = Pt(0.5)
        tf = rb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val; r.font.size = Pt(10)
        r.font.color.rgb = txt_c if j > 0 else DGRAY
        if i == 0: r.font.bold = True

box(sl4, 0.35, 7.12, 12.6, 0.25,
    "MLP+EN Long-Short outperforms Buy & Hold by +32 pp  ·  "
    "All strategies modeled with 0.05 % per-trade transaction costs",
    font_size=9, color=RGBColor(0x75,0x75,0x75), align=PP_ALIGN.CENTER)


out = "VIP_Presentation_Option3.pptx"
prs.save(out)
print(f"Saved: {out}")
